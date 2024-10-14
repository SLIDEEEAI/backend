import hashlib

import os, sys

import json

from json import dumps
from re import findall
from io import BytesIO

from django.conf import settings
from django.core.files.base import File

from openai.types.chat import ChatCompletion

from requests import get as request_get

from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE

from datetime import datetime

from .models import Picture


def chat_competions_create(system_content: str) -> str | None:
    # print(system_content)
    # sys.stdout.flush()
    chat_completion: ChatCompletion = settings.OPENAI_CLIENT.chat.completions.create(
        messages=[
            {
                "role": "system",
                "content": system_content,
            }
        ],
        model="gpt-3.5-turbo",
    )

    if chat_completion.choices:
        return chat_completion.choices[0].message.content
    else:
        return None


def images_generate(prompt: str):
    response = settings.OPENAI_CLIENT.images.generate(
        model="dall-e-2",
        prompt=prompt,
        size="1024x1024",
        quality="standard",
        n=1,
    )

    return response.data[0].url if response.data else None



def generate_images_from_list(prompts: list[str]):
    return [
        images_generate(prompt) for prompt in prompts
    ]

import hashlib
import openai
def generate_images2(presentation_theme, num_images):
    try:
        # Генерация изображений с помощью OpenAI API
        response = openai.Image.create(
            prompt=presentation_theme,
            n=num_images,
            size="1024x1024",
            quality="hd",
        )

        images = response['data']
        saved_images = []

        for image_data in images:
            image_url = image_data['url']
            # Загрузка изображения из URL
            image_content = requests.get(image_url).content
            # Создание хэш-имени для изображения
            hash_name = hashlib.sha256(image_content).hexdigest()
            extension = "png"  # или другой подходящий формат
            picture_name = f'{hash_name}.{extension}'

            # Проверка наличия изображения в базе данных
            picture = Picture.objects.filter(hash_name=picture_name).first()
            if not picture:
                picture = Picture(hash_name=picture_name)
                picture_file = ContentFile(image_content)
                picture.source.save(picture_name, picture_file)
                picture.save()

            saved_images.append({
                "url": picture.source.url,
                "hash_name": hash_name,
            })

        return saved_images

    except Exception  as e:
        print(f"Error generating images: {e}")
        return []


def generate_slides_theme(presentation_theme: str, slides_count: int) -> list[str]:
    content = chat_competions_create(f"Write me short topics for slide headings in the amount of {slides_count} pieces. for a presentation on the topic '{presentation_theme}'' in Russian, in the form of a list without numbering, any enumeration, and any additional words.")
    # content = chat_competions_create(f"Write me slide themes up to 10 words each in the amount of {slides_count} pcs. for a presentation on the topic {presentation_theme} in Russian, in a list without numbering or any enumeration, as well as any additional words!")

    if content:
        for x in findall(r"[^\W]\w+.*", content)[:slides_count]:
            yield x.strip()
    else:
        yield None


def generate_slides_text(slides_themes: list) -> list[str]:
    system_content = "You should write several sentences on each specified topic in Russian - {0}. Sentences should not exceed 500 words and each should contain a complete thought. You can also choose not to add words from yourself, briefly and only on the topic! Your complete answer should contain exactly - {1} sentences. The topics themselves should not appear in your answer!"
    content = chat_competions_create(system_content.format(", ".join(slides_themes), len(slides_themes)))

    print(f'{content=}')
    sys.stdout.flush()

    if content:
        for x in findall(r"\w+.*", content)[:len(slides_themes)]:
            yield x
    else:
        yield None


################### НОВЫЕ АПИ

def generate_custom_request(prompt: str, max_tokens: int = 500) -> str:
    # Генерация короткого текста для введения или краткого описания темы презентации.
    messages = [
        {"role": "user", "content": prompt}
    ]
    response = settings.OPENAI_CLIENT.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=messages,
        max_tokens=max_tokens,
    )
    return response.choices[0].message.content.strip()


def generate_short_text(presentation_theme: str, max_tokens: int = 50) -> str:
    # Генерация короткого текста для введения или краткого описания темы презентации.
    messages = [
        {"role": "system", "content": "You are a helpful assistant that generates short introductions or brief descriptions for presentation themes."},
        {"role": "user", "content": f"Напишите короткое введение или краткое описание для темы презентации: {presentation_theme}. Ограничьте свой ответ {max_tokens} токенами."}
    ]
    response = settings.OPENAI_CLIENT.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=messages,
        max_tokens=max_tokens,
    )
    return response.choices[0].message.content.strip()


def images_generate_4(prompt: str):
    response = settings.OPENAI_CLIENT.images.generate(
        model="dall-e-3",
        prompt=prompt,
        size="1024x1024",
        quality="standard",
        n=1,
    )

    return response.data[0].url if response.data else None

def generate_long_text(presentation_theme: str, max_tokens: int = 50) -> str:
    # Генерация длинного текста для подробного описания содержания презентации или конкретного слайда.
    messages = [
        {"role": "system", "content": "You should write several sentences on each specified topic in Russian - {0}. Sentences should not exceed 500 words and each should contain a complete thought. You can also choose not to add words from yourself, briefly and only on the topic! Your complete answer should contain exactly - {1} sentences. The topics themselves should not appear in your answer!"},
        {"role": "user", "content": f"Напишите про {presentation_theme}. Пиши только текст, не надо дополнительных вводных. Мне надо это будет скопировать на слайд. Не пиши название самой темы. Ваш ответ не должен превышать {max_tokens} токенов."}
    ]
    response = settings.OPENAI_CLIENT.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=messages,
        max_tokens=max_tokens,
    )
    return response.choices[0].message.content.strip()



def generate_bullet_points(presentation_theme: str, max_items: int = 5) -> list[str]:
    # Генерация списка маркированных или нумерованных пунктов, который может использоваться для перечисления ключевых моментов или идей.
    messages = [
        {"role": "system", "content": "You are a helpful assistant that generates bullet points related to presentation themes."},
        {"role": "user", "content": f"Напишите функции и виды по {max_items} пунктов, связанных с темой презентации: {presentation_theme}."}
    ]
    response = settings.OPENAI_CLIENT.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=messages,
        max_tokens=100,  # Предположим, что каждый пункт составляет около 20 токенов
    )
    items = response.choices[0].message.content.split("\n")[:max_items]
    return [item.strip() for item in items]



def generate_image_with_caption(presentation_theme: str) -> tuple[str, str]:
    # Генерация изображения и соответствующей подписи, которые могут использоваться для иллюстрации концепций или примеров в презентации.
    prompt = f"Сгенерируйте изображение и подпись, связанные с темой презентации: {presentation_theme}."
    response = settings.OPENAI_CLIENT.images.generate(
        model="dall-e-3",
        prompt=prompt,
        size="1024x1024",
        quality="hd",
        n=1,
    )
    if response.data:
        image_url = response.data[0].url
        caption = generate_short_text(presentation_theme)
        return image_url, caption
    else:
        return None, None


def generate_quote(presentation_theme: str) -> str:
    # Генерация цитаты или высказывания.
    messages = [
        {"role": "system", "content": "You are a helpful assistant that generates quotes or sayings related to presentation themes."},
        {"role": "user", "content": f"Напишите цитату или высказывание, связанные с темой презентации: {presentation_theme}."}
    ]
    response = settings.OPENAI_CLIENT.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=messages,
        max_tokens=50,  # цитата короткая - 100 слов
    )
    return response.choices[0].message.content.strip()



def generate_chart_data(presentation_theme: str) -> str:
    # Генерация данных для построения графиков или диаграмм, которые могут визуализировать статистику или тренды.
    messages = [
        {"role": "system", "content": "You are a helpful assistant that generates data for creating charts or diagrams related to presentation themes."},
        {"role": "user", "content": f"Сгенерируйте данные для создания графиков или диаграмм, связанных с темой презентации: {presentation_theme}."}
    ]
    response = settings.OPENAI_CLIENT.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=messages,
        max_tokens=100,  # Предположим, что короткий ответ на данные для графика
    )
    return response.choices[0].message.content.strip()

def generate_questions(presentation_theme: str, max_questions: int = 3) -> list[str]:
    # Генерация вопросов, которые могут использоваться для вовлечения аудитории или обсуждения темы презентации.
    messages = [
        {"role": "system", "content": "You are a helpful assistant that generates questions related to presentation themes for audience engagement or discussion."},
        {"role": "user", "content": f"Сгенерируйте {max_questions} вопроса, связанных с темой презентации: {presentation_theme}."}
    ]
    response = settings.OPENAI_CLIENT.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=messages,
        max_tokens=150,  # Предположим, что каждый вопрос около 50 токенов
    )
    questions = response.choices[0].message.content.split("\n")[:max_questions]
    return [question.strip() for question in questions]

def generate_slide_title(presentation_theme: str) -> str:
    # Генерация заголовка слайда, который может четко описывать содержание следующего слайда.
    messages = [
        {"role": "system", "content": "You are a helpful assistant that generates slide titles that clearly describe the content related to the presentation theme."},
        {"role": "user", "content": f"Сгенерируйте заголовок слайда, который четко описывает содержание, связанное с темой презентации: {presentation_theme}."}
    ]
    response = settings.OPENAI_CLIENT.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=messages,
        max_tokens=20,  # Предположим, что заголовок короткий
    )
    return response.choices[0].message.content.strip()

def generate_slide_heading(presentation_theme: str) -> str:
    # Генерация заголовка раздела презентации, который может помочь структурировать ее содержание.
    messages = [
        {"role": "system", "content": "You are a helpful assistant that generates slide headings that help structure the content related to the presentation theme."},
        {"role": "user", "content": f"Сгенерируйте заголовок раздела презентации, который поможет структурировать ее содержание, связанное с темой: {presentation_theme}."}
    ]
    response = settings.OPENAI_CLIENT.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=messages,
        max_tokens=30,  # Предположим, что заголовок немного длиннее
    )
    return response.choices[0].message.content.strip()

def generate_images(presentation_theme: str, num_images: int = 3) -> list[str]:
    # Генерация изображений, связанных с темой презентации.
    images = []

    for _ in range(num_images):
        prompt = f"Сгенерируйте изображение по теме презентации: {presentation_theme}."
        response = settings.OPENAI_CLIENT.images.generate(
            model="dall-e-2",
            prompt=prompt,
            size="1024x1024",
            quality="hd",
            n=1,
        )
        if response.data:
            images.append(response.data[0].url)

    return images


def pt_to_str(pt_value):
    """КОНВЕРТАЦИЯ ПТ В ПИСКЕЛИ"""
    if pt_value is not None:
        return f"{int(pt_value.pt)}px"
    return None


def add_empty_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def set_slide_dimensions(prs, width, height):
    prs.slide_width = Inches(width)
    prs.slide_height = Inches(height)


def set_slide_background(slide, rgb_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb_color)


def add_text_box(slide, left, top, width, height, text, font_size, is_bold=False):
    text_frame = slide.shapes.add_textbox(left, top, width, height).text_frame
    text_frame.word_wrap = True
    text_frame.auto_size = True
    paragraph = text_frame.add_paragraph()
    paragraph.text = text
    paragraph.font.bold = is_bold
    paragraph.font.size = Pt(font_size)
    paragraph.font.name = 'Arial'


def add_picture(slide, img_path, left, top, width, height):
    slide.shapes.add_picture(
        img_path, left=left, top=top, width=width, height=height
    )


def get_bytesio_from_url(url):
    return BytesIO(request_get(url).content)


def generate_pptx_from_json(slides_data):
    presentation_object = Presentation()

    for slide_data in slides_data:
        slide = add_empty_slide(presentation_object)

        # Set slide dimensions and background color
        set_slide_dimensions(presentation_object, 26.67, 15)  # Assuming the same dimensions for all slides
        set_slide_background(slide, (255, 248, 220))  # Assuming the same background color for all slides

        # Add text boxes and pictures to the slide
        add_text_box(slide, Pt(60), Pt(40), Pt(757), Pt(300), slide_data['title'], 110, is_bold=True)
        # add_text_box(slide, Pt(60), Pt(321), Pt(245), Pt(45), slide_data['main_text'], 60)
        add_text_box(slide, Pt(60), Pt(450), Pt(757), Pt(430), slide_data['text'], 35)

        add_picture(slide, get_bytesio_from_url(slide_data['image']), Pt(957), Pt(80), Pt(883), Pt(920))

    return presentation_object


def generate_json_from_pptx(presentation_object, title: str = 'Презентация',):
    try:
        json_object = {
            # 'id': 0, # Вынесено
            # 'author': None, # Вынесено
            'group': None,
            "favourite" : False,
            "removed" : False,
            "date_created" : int(datetime.now().timestamp()),
            "date_edited" : int(datetime.now().timestamp()),
            "theme" : {
                'background_color': presentation_object.slides[0].background.fill.fore_color.rgb,
                'font_info': {
                    "titles" : {
                        "name" : 'Calibri',
                        "size" : 44,
                        "bold" : True,
                        "italic" : False
                    },
                    "main_texts" : { 
                        "name" : 'Calibri',
                        "size" : 18,
                        "bold" : False,
                        "italic" : False
                    }
                }
            },
            "len_slides" : 0,
            'title': title,
            'slides': []
        }

        for idx, slide in enumerate(presentation_object.slides):
            slide_info = {
                'index': idx,
                'theme' : {
                   'background_color' : None,                       
                   'background_img' : None                 
                },
                'content': {
                    'textboxes': [],
                    "pictures" : [],
                    "tables" : [],
                    "figures" : [],
                    "graphics" : []
                }
            }

            textbox_counter = 0

            # Временное решение для заголовка
            for counter, shape in enumerate(slide.shapes):
                # Временное решение для заголовка
                if shape.has_text_frame:
                    # Заголовок
                    if textbox_counter == 0: 
                        slide_info['content']['textboxes'].append({
                            'type': 'h2',
                            'left': pt_to_str(shape.left),
                            'top': pt_to_str(shape.top),
                            'width': pt_to_str(shape.width),
                            'height': pt_to_str(shape.height),
                            'text': shape.text,
                            'font_size': pt_to_str(shape.text_frame.paragraphs[0].font.size),
                            'bold': shape.text_frame.paragraphs[0].font.bold,
                            'background_color': None,
                            'italic': False,
                            'align': None
                        })
                    # Доп. текст(описание)
                    else:
                        slide_info['content']['textboxes'].append({
                            'type': 'p',
                            'left': pt_to_str(shape.left),
                            'top': pt_to_str(shape.top),
                            'width': pt_to_str(shape.width),
                            'height': pt_to_str(shape.height),
                            'text': shape.text,
                            'font_size': pt_to_str(shape.text_frame.paragraphs[0].font.size),
                            'bold': shape.text_frame.paragraphs[0].font.bold,
                            'background_color': None,
                            'italic': False,
                            'align': None
                        })
                    textbox_counter += 1
                elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # Picture shape type
                    image_blob = str(shape.image.blob)
                    extension = shape.image.ext

                    hash_name = hashlib.sha256(image_blob.encode()).hexdigest()
                    picture_name = f'{hash_name}.{extension}'

                    picture = Picture.objects.filter(hash_name=picture_name).first()
                    if not picture:
                        picture = Picture.objects.create(hash_name=picture_name)

                    if not picture.source or not bool(picture.source):
                        with BytesIO(shape.image.blob) as stream:
                            picture_file = File(stream)
                            picture.source.save(picture_name, picture_file)
                        picture.save()

                    slide_info['content']['pictures'].append({
                        'left': pt_to_str(shape.left),
                        'top': pt_to_str(shape.top),
                        'width': pt_to_str(shape.width),
                        'height': pt_to_str(shape.height),
                        "url": picture.source.url,
                        'background_color': None
                    })
               

            json_object['slides'].append(slide_info)
            json_object['len_slides'] = len(json_object['slides'])
    except Exception as exc:
        print('ERROR : ', exc)
        return {}

    return json_object


def generate_json_object(themes: list[str], slides: list[str], images: list[str]) -> str:
    
    # eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJ0b2tlbl90eXBlIjoiYWNjZXNzIiwiZXhwIjoxNzEyMTEwOTMxLCJpYXQiOjE3MTIwMjQ1MzEsImp0aSI6ImM5NjA4YmNkYTQzYzRjNGZiM2QyMmIyZTViZDU4MjA4IiwidXNlcl9pZCI6NSwidXNlcm5hbWUiOiJJbHlhIiwiZW1haWwiOiJ0ZXRzQGJrLnJ1In0.GlyddsxaZsrXj-iUgJfxpXsbl-YdJUJOelGmLWe24pg

    print(themes)
    print(slides)

    return dumps(
        generate_json_from_pptx(
            generate_pptx_from_json([
                {
                    'title': themes[i],
                    'text': slide,
                    'image': images[i]
                } for i, slide in enumerate(slides)
            ])
        )
    )


def export_presentation(presentation, presentation_type) -> str:
    presentation_object = json.loads(presentation.json)

    pptx_object = generate_pptx_from_json([
        {
            'title': slide['content']['textboxes'][0]['text'],
            'text': slide['content']['textboxes'][1]['text'],
            'image': settings.BASE_URL + slide['content']['pictures'][0]['url']
        } for i, slide in enumerate(presentation_object['slides'])
    ])

    # Если нужно отдать pdf, то генерируем pdf из pptx
    if presentation_type == 'pdf':
        # Преобразование pdf из pptx
        # ЗДЕСЬ КОД ПО ПРЕОБРАЗОВАНИЮ

        # Нужно присвоить обьект
        presentation_object = None
    else:
        presentation_object = pptx_object

    name = f'/{presentation.id}-{presentation.user.id}-{int(datetime.now().timestamp())}.pptx'

    os_path = str(settings.MEDIA_ROOT) + '/pptx' + name

    # Сохраняем презу и отдаем путь ПОЛНЫЙ(ссылка)
    presentation_object.save(os_path)

    url_path = settings.BASE_URL + '/media/pptx' + name

    return url_path